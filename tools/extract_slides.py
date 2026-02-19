import json
import os
from pathlib import Path

import pdfplumber
from pptx import Presentation


def extract_pptx_text(path: Path) -> list[dict]:
    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        slides.append({"slide": idx, "text": texts})
    return slides


def extract_pdf_text(path: Path) -> list[dict]:
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            pages.append({"page": idx, "text": lines})
    return pages


def main() -> None:
    base_dir = Path(r"c:\VScode\4IE2\.tex\rinri\data")
    out_path = base_dir / "_extracted_slides.json"
    data = {}

    for path in sorted(base_dir.glob("*.pptx")):
        data[path.name] = {
            "type": "pptx",
            "slides": extract_pptx_text(path),
        }

    for path in sorted(base_dir.glob("*.pdf")):
        data[path.name] = {
            "type": "pdf",
            "pages": extract_pdf_text(path),
        }

    out_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    print(str(out_path))


if __name__ == "__main__":
    main()
